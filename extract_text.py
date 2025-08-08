#!/usr/bin/env python3
"""
Document text extraction module for PDF and PPTX files.

Uses PyMuPDF for PDFs with OCR fallback via pytesseract.
Uses python-pptx for PPTX files.
"""

import io
import json
import logging
import os
import sys
from pathlib import Path
from typing import Dict, List
from uuid import uuid4

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None


def extract_document_text(file_path: str) -> Dict:
    """
    Extracts all text from a PDF or PPTX file.
    Uses PyMuPDF for PDFs, with OCR fallback via pytesseract if a page has no text.
    Uses python-pptx for PPTX.
    
    Args:
        file_path: Path to the PDF or PPTX file
        
    Returns:
        A dict with shape:
        {
          "doc": { "title": <basename of file>, "pages": <int total pages/slides> },
          "sections": [
            { "id": "s1", "title": <section or slide title or "Page N">, "text": <extracted text> },
            ...
          ]
        }
        
    Raises:
        ValueError: If file type is unsupported or required dependencies are missing
        FileNotFoundError: If file doesn't exist
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_path_obj = Path(file_path)
    file_extension = file_path_obj.suffix.lower()
    file_title = file_path_obj.stem
    
    if file_extension == '.pdf':
        return _extract_pdf_text(file_path, file_title)
    elif file_extension == '.pptx':
        return _extract_pptx_text(file_path, file_title)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}. Only .pdf and .pptx are supported.")


def _extract_pdf_text(file_path: str, file_title: str) -> Dict:
    """Extract text from PDF using PyMuPDF with OCR fallback."""
    if fitz is None:
        raise ValueError("PyMuPDF (fitz) is required for PDF processing. Install with: pip install PyMuPDF")
    
    doc = fitz.open(file_path)
    sections = []
    
    for page_num in range(len(doc)):
        try:
            page = doc.load_page(page_num)
            text = page.get_text("text").strip()
            
            # If no text found, try OCR
            if not text:
                text = _ocr_page(page, page_num + 1)
            
            if text:  # Only add non-empty sections
                section = {
                    "id": uuid4().hex[:8],
                    "title": f"Page {page_num + 1}",
                    "text": text
                }
                sections.append(section)
                
        except Exception as e:
            logging.error(f"Error processing page {page_num + 1}: {e}")
            continue
    
    doc.close()
    
    return {
        "doc": {
            "title": file_title,
            "pages": len(doc)
        },
        "sections": sections
    }


def _extract_pptx_text(file_path: str, file_title: str) -> Dict:
    """Extract text from PPTX using python-pptx."""
    if Presentation is None:
        raise ValueError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
    
    presentation = Presentation(file_path)
    sections = []
    
    for slide_num, slide in enumerate(presentation.slides, 1):
        try:
            # Extract title
            title = _get_slide_title(slide, slide_num)
            
            # Extract all text from shapes
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            
            # Combine all text
            text = "\n".join(text_parts).strip()
            
            if text:  # Only add non-empty sections
                section = {
                    "id": uuid4().hex[:8],
                    "title": title,
                    "text": text
                }
                sections.append(section)
                
        except Exception as e:
            logging.error(f"Error processing slide {slide_num}: {e}")
            continue
    
    return {
        "doc": {
            "title": file_title,
            "pages": len(presentation.slides)
        },
        "sections": sections
    }


def _get_slide_title(slide, slide_num: int) -> str:
    """Extract title from slide or return default."""
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    
    return f"Slide {slide_num}"


def _ocr_page(page, page_num: int) -> str:
    """Perform OCR on a PDF page."""
    if pytesseract is None or Image is None:
        logging.warning(f"Page {page_num}: OCR dependencies not available (pytesseract, Pillow)")
        return ""
    
    try:
        # Render page as image at ~200 DPI
        mat = fitz.Matrix(2.0, 2.0)  # 2.0 scale ≈ 200 DPI for typical 72 DPI PDF
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        
        # Convert to PIL Image and OCR
        img = Image.open(io.BytesIO(img_data))
        text = pytesseract.image_to_string(img).strip()
        
        if text:
            print(f"Used OCR for Page {page_num}")
        
        return text
        
    except Exception as e:
        logging.error(f"OCR failed for page {page_num}: {e}")
        return ""


def main():
    """CLI interface for testing."""
    if len(sys.argv) != 2:
        print("Usage: python extract_text.py <file_path>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    try:
        result = extract_document_text(file_path)
        print(json.dumps(result, indent=2, ensure_ascii=False))
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    # Configure logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(levelname)s: %(message)s'
    )
    
    main()
